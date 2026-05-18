"""
AI Income Lab — PowerPoint v2
Clean, professional dark theme. Properly structured slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colours ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0A, 0x0A)
CARD     = RGBColor(0x1A, 0x1A, 0x1A)
CARD2    = RGBColor(0x22, 0x22, 0x22)
LIME     = RGBColor(0xC4, 0xFF, 0x3D)
ORANGE   = RGBColor(0xFF, 0x78, 0x49)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xA8, 0xA8, 0x9E)
DGRAY    = RGBColor(0x5A, 0x5A, 0x52)
BLACK    = RGBColor(0x00, 0x00, 0x00)

W = Inches(13.3)
H = Inches(7.5)

# ── Low-level helpers ─────────────────────────────────────────────────────────

def hex3(c): return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def spPr(shape):
    el = shape._element
    sp = el.find(qn('p:spPr'))
    if sp is None:
        sp = etree.SubElement(el, qn('p:spPr'))
    return sp

def fill_solid(shape, color, alpha=100):
    pr = spPr(shape)
    for t in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill')]:
        e = pr.find(t); e is not None and pr.remove(e)
    sf = etree.SubElement(pr, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', hex3(color))
    if alpha < 100:
        etree.SubElement(sc, qn('a:alpha')).set('val', str(alpha * 1000))

def fill_none(shape):
    pr = spPr(shape)
    for t in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill')]:
        e = pr.find(t); e is not None and pr.remove(e)
    etree.SubElement(pr, qn('a:noFill'))

def line_solid(shape, color, pt=1, alpha=100):
    pr = spPr(shape)
    ln = pr.find(qn('a:ln'))
    if ln is not None: pr.remove(ln)
    ln = etree.SubElement(pr, qn('a:ln'))
    ln.set('w', str(int(pt * 12700)))
    sf = etree.SubElement(ln, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', hex3(color))
    if alpha < 100:
        etree.SubElement(sc, qn('a:alpha')).set('val', str(alpha * 1000))

def line_none(shape):
    pr = spPr(shape)
    ln = pr.find(qn('a:ln'))
    if ln is not None: pr.remove(ln)
    ln = etree.SubElement(pr, qn('a:ln'))
    etree.SubElement(ln, qn('a:noFill'))

def add_shadow(shape):
    pr = spPr(shape)
    el = pr.find(qn('a:effectLst'))
    if el is None: el = etree.SubElement(pr, qn('a:effectLst'))
    s = etree.SubElement(el, qn('a:outerShdw'))
    s.set('blurRad', str(int(20 * 12700)))
    s.set('dist',    str(int(6  * 12700)))
    s.set('dir', str(int(135 * 60000)))
    s.set('algn', 'ctr'); s.set('rotWithShape', '0')
    sc = etree.SubElement(s, qn('a:srgbClr'))
    sc.set('val', '000000')
    etree.SubElement(sc, qn('a:alpha')).set('val', '65000')

def fill_grad_v(shape, stops):
    """stops = [(pos 0-100, RGBColor), ...]"""
    pr = spPr(shape)
    for t in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill')]:
        e = pr.find(t); e is not None and pr.remove(e)
    gf = etree.SubElement(pr, qn('a:gradFill'))
    gl = etree.SubElement(gf, qn('a:gsLst'))
    for pos, col in stops:
        gs = etree.SubElement(gl, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        sc = etree.SubElement(gs, qn('a:srgbClr'))
        sc.set('val', hex3(col))
    lin = etree.SubElement(gf, qn('a:lin'))
    lin.set('ang', '5400000'); lin.set('scaled', '0')

# ── Shape/text builders ───────────────────────────────────────────────────────

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl

def box(sl, x, y, w, h, fill=CARD, border=None, border_pt=1, shadow=False):
    s = sl.shapes.add_shape(1, x, y, w, h)
    fill_solid(s, fill)
    if border:
        line_solid(s, border, border_pt)
    else:
        line_none(s)
    if shadow:
        add_shadow(s)
    return s

def accent_bar(sl, x, y, h, color=LIME):
    s = sl.shapes.add_shape(1, x, y, Pt(4), h)
    fill_solid(s, color)
    line_none(s)

def hbar(sl, x, y, w, color=DGRAY, pt=1):
    s = sl.shapes.add_shape(1, x, y, w, Pt(pt))
    fill_solid(s, color)
    line_none(s)

def txt(sl, text, x, y, w, h,
        size=14, color=WHITE, font='Calibri',
        bold=False, italic=False, align=PP_ALIGN.LEFT,
        wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb

def mtxt(sl, lines, x, y, w, h,
         size=13, color=GRAY, font='Calibri', spacing=5):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(spacing)
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color

def arrow_list(sl, items, x, y, w, h, size=12):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        r1 = p.add_run(); r1.text = '→  '
        r1.font.name = 'Courier New'; r1.font.size = Pt(size); r1.font.color.rgb = LIME
        r2 = p.add_run(); r2.text = item
        r2.font.name = 'Calibri'; r2.font.size = Pt(size); r2.font.color.rgb = WHITE

def headline(sl, part1, part2, x, y, w, size=38):
    tb = sl.shapes.add_textbox(x, y, w, Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = part1
    r1.font.name = 'Georgia'; r1.font.size = Pt(size); r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = part2
    r2.font.name = 'Georgia'; r2.font.size = Pt(size)
    r2.font.color.rgb = GRAY; r2.font.italic = True

# ── Section header shared element ─────────────────────────────────────────────

def section_header(sl, label, h1, h2, sub):
    # left label column
    txt(sl, label, Inches(0.45), Inches(0.38), Inches(1.9), Pt(18),
        size=9, color=LIME, font='Courier New', bold=True)
    hbar(sl, Inches(0.45), Inches(0.62), Inches(1.9), LIME, 1)
    # headline
    headline(sl, h1, h2, Inches(2.6), Inches(0.18), Inches(10.3))
    # sub
    txt(sl, sub, Inches(2.6), Inches(1.25), Inches(10.3), Pt(36),
        size=13, color=GRAY)
    hbar(sl, Inches(0.45), Inches(2.05), Inches(12.4))


# ════════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    sl = new_slide(prs)

    # subtle vertical grid lines
    for i in range(14):
        s = sl.shapes.add_shape(1, Inches(i), Inches(0), Pt(1), H)
        fill_solid(s, RGBColor(0x22,0x22,0x22))
        line_none(s)

    # logo block
    lb = sl.shapes.add_shape(1, Inches(0.45), Inches(0.45), Pt(38), Pt(38))
    fill_solid(lb, LIME); line_none(lb)
    txt(sl, 'AI', Inches(0.45), Inches(0.45), Pt(38), Pt(38),
        size=14, color=BLACK, font='Courier New', bold=True, align=PP_ALIGN.CENTER)
    txt(sl, 'INCOME LAB / personal notebook',
        Inches(1.0), Inches(0.55), Inches(5), Pt(22),
        size=11, color=WHITE, font='Courier New', bold=True)

    # meta
    txt(sl, '// Field manual  ·  Edition 01 — May 2026',
        Inches(0.45), Inches(1.45), Inches(8), Pt(20),
        size=10, color=DGRAY, font='Courier New')

    # live tag pill
    pill = sl.shapes.add_shape(9, Inches(0.45), Inches(2.05), Inches(3.2), Pt(26))
    fill_solid(pill, RGBColor(0x14,0x28,0x00))
    line_solid(pill, LIME, 1)
    txt(sl, '● Live transmission — AI economy 2026',
        Inches(0.65), Inches(2.07), Inches(3.0), Pt(22),
        size=9.5, color=LIME, font='Courier New')

    # main headline
    tb = sl.shapes.add_textbox(Inches(0.45), Inches(2.65), Inches(10.2), Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'Making money online,\npowered by '
    r1.font.name = 'Georgia'; r1.font.size = Pt(58); r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = 'artificial intelligence'
    r2.font.name = 'Georgia'; r2.font.size = Pt(58)
    r2.font.color.rgb = LIME; r2.font.italic = True

    txt(sl,
        'A working notebook of the methods, platforms, and tools that actually generate income '
        'in 2026 — researched, tested, and documented for personal study. No hype. Just leverage.',
        Inches(0.45), Inches(5.75), Inches(9.5), Pt(52),
        size=15, color=GRAY, wrap=True)

    # right info box
    bx = box(sl, Inches(10.6), Inches(2.75), Inches(2.3), Inches(2.0), CARD2, RGBColor(0x33,0x33,0x33), shadow=True)
    txt(sl, '// operator',   Inches(10.75), Inches(2.9),  Inches(2.0), Pt(16), size=9,  color=DGRAY, font='Courier New')
    txt(sl, 'Personal\nstudy log', Inches(10.75), Inches(3.12), Inches(2.0), Inches(0.9), size=19, color=WHITE, font='Georgia')
    txt(sl, '// edition',    Inches(10.75), Inches(3.95), Inches(2.0), Pt(16), size=9,  color=DGRAY, font='Courier New')
    txt(sl, 'May 2026',      Inches(10.75), Inches(4.17), Inches(2.0), Pt(22), size=14, color=LIME, font='Courier New')


def s02_market(prs):
    sl = new_slide(prs)
    section_header(sl,
        '01 / Market',
        'The opportunity, ',
        'by the numbers.',
        'Why now matters. The gap between tools and those who can wield them is the largest it has ever been.')

    stats = [
        ('674', 'B',  'Global gig economy 2026',       '$ — global'),
        ('70',  '%',  'Firms increasing AI spend',      '% — companies'),
        ('15.7','T',  "AI GDP contribution by 2030",    '$ — projection'),
        ('300', '%',  'AI job listings since 2023',     '% — growth'),
    ]

    sw, sh = Inches(3.06), Inches(2.35)
    gap = Inches(0.115)
    sy = Inches(2.25)
    for i, (num, unit, desc, tag) in enumerate(stats):
        x = Inches(0.42) + i * (sw + gap)
        c = box(sl, x, sy, sw, sh, CARD, RGBColor(0x2a,0x2a,0x2a), shadow=True)
        txt(sl, tag,  x+Inches(0.2), sy+Pt(12), sw-Inches(0.4), Pt(18),
            size=9, color=DGRAY, font='Courier New')
        # big number
        nb = sl.shapes.add_textbox(x+Inches(0.2), sy+Pt(30), sw-Inches(0.4), Inches(1.05))
        p = nb.text_frame.paragraphs[0]
        r1 = p.add_run(); r1.text = num
        r1.font.name = 'Georgia'; r1.font.size = Pt(62); r1.font.color.rgb = LIME
        r2 = p.add_run(); r2.text = unit
        r2.font.name = 'Georgia'; r2.font.size = Pt(30); r2.font.color.rgb = GRAY
        txt(sl, desc, x+Inches(0.2), sy+Inches(1.55), sw-Inches(0.4), Pt(40),
            size=11, color=GRAY, font='Courier New')

    txt(sl,
        'AI consulting: $100–$300/hr  ·  Prompt engineering: $50–$150/hr  ·  Most digital products take days, not weeks',
        Inches(0.42), Inches(6.82), Inches(12.4), Pt(22),
        size=10, color=DGRAY, font='Courier New', align=PP_ALIGN.CENTER)


def s03_methods_a(prs):
    sl = new_slide(prs)
    section_header(sl,
        '02 / Methods  (1–5)',
        'Ten ways to earn, ',
        'ranked by leverage.',
        'Every method pairs a real human skill with an AI tool. The goal: deliver outcomes businesses already pay for, faster.')

    methods = [
        ('001', 'AI-powered\ncontent services',
         'Ghost-write blogs, newsletters & landing copy. AI handles research and first drafts. You handle voice & strategy.',
         'Beginner', '$40–$100/hr', 'ChatGPT · Surfer SEO'),
        ('002', 'AI art &\nprint-on-demand',
         'Generate visuals with Midjourney, sell as prints & merch on Etsy / Redbubble. Pick a narrow niche.',
         'Beginner', '$500–$3K/mo', 'Midjourney · Printify'),
        ('003', 'Templates,\nprompt packs & eBooks',
         'Build once, sell forever. Notion templates, prompt libraries, resume kits. Clearest path to passive income.',
         'Beginner', 'Semi-passive', 'Notion · Gumroad'),
        ('004', 'Faceless\nYouTube channels',
         'Scripted with AI, narrated with ElevenLabs, edited with Pictory. Long ramp, compounds powerfully.',
         'Intermediate', '$150–$2K/mo', 'ElevenLabs · Pictory'),
        ('005', 'Voiceover &\nlocalization',
         'Sell narration, dubbing & translation for video, audiobooks & ads. High-margin, AI completes in minutes.',
         'Intermediate', '$50–$500/proj', 'ElevenLabs · Murf'),
    ]

    mw, mh = Inches(2.54), Inches(2.88)
    gap = Inches(0.07)
    my = Inches(2.2)
    for i, (num, name, desc, lvl, pay, stack) in enumerate(methods):
        x = Inches(0.42) + i * (mw + gap)
        c = box(sl, x, my, mw, mh, CARD, RGBColor(0x2a,0x2a,0x2a), shadow=True)
        accent_bar(sl, x, my, mh)
        txt(sl, num, x+Inches(0.22), my+Pt(10), mw, Pt(18),
            size=9, color=DGRAY, font='Courier New')
        txt(sl, name, x+Inches(0.22), my+Pt(26), mw-Inches(0.3), Pt(54),
            size=17, color=WHITE, font='Georgia', wrap=True)
        txt(sl, desc, x+Inches(0.22), my+Pt(92), mw-Inches(0.3), Pt(74),
            size=11, color=GRAY, wrap=True)
        # level chip
        lv = sl.shapes.add_shape(1, x+Inches(0.22), my+Inches(2.28), Inches(1.02), Pt(21))
        fill_solid(lv, RGBColor(0x1a,0x28,0x00)); line_solid(lv, RGBColor(0x4a,0x70,0x00), 1)
        txt(sl, lvl, x+Inches(0.24), my+Inches(2.29), Inches(1.0), Pt(19),
            size=9, color=LIME, font='Courier New')
        pv = sl.shapes.add_shape(1, x+Inches(1.30), my+Inches(2.28), Inches(1.1), Pt(21))
        fill_solid(pv, CARD2); line_solid(pv, RGBColor(0x33,0x33,0x33), 1)
        txt(sl, pay, x+Inches(1.32), my+Inches(2.29), Inches(1.08), Pt(19),
            size=9, color=GRAY, font='Courier New')
        hbar(sl, x+Inches(0.22), my+Inches(2.57), mw-Inches(0.3), DGRAY)
        txt(sl, stack, x+Inches(0.22), my+Inches(2.63), mw-Inches(0.3), Pt(18),
            size=9, color=DGRAY, font='Courier New')


def s04_methods_b(prs):
    sl = new_slide(prs)
    section_header(sl,
        '02 / Methods  (6–10)',
        'Advanced methods, ',
        'higher leverage.',
        'These take longer to ramp but deliver the highest hourly rates and the most durable income.')

    methods = [
        ('006', 'Affiliate marketing\nwith AI SEO',
         'Build niche review sites or YouTube. Use AI for keywords, content & B-roll.',
         'Intermediate', '20–50% comm.', 'Surfer SEO · Ahrefs'),
        ('007', 'Online courses\n& cohorts',
         'Teach one skill exceptionally well. AI compresses design, slides & marketing. Highest margin product.',
         'Intermediate', '$1K–$10K/course', 'Thinkific · Notion'),
        ('008', 'Custom GPTs\n& AI agents',
         'Narrow-purpose bots for small businesses — support, lead-qualify, knowledge bases.',
         'Advanced', '$75–$200/hr', 'OpenAI · Voiceflow'),
        ('009', 'Business workflow\nautomation',
         'Wire AI into operations. Email triage, meeting notes, lead routing. 4-figure retainers.',
         'Advanced', '$75–$200/hr', 'Make · n8n · Zapier'),
        ('010', 'AI consulting\n& training',
         'Help orgs find where AI creates real ROI. Run workshops. Audit workflows. Highest hourly rate.',
         'Advanced', '$100–$300/hr', 'Your expertise + Loom'),
    ]

    mw, mh = Inches(2.54), Inches(2.88)
    gap = Inches(0.07)
    my = Inches(2.2)
    for i, (num, name, desc, lvl, pay, stack) in enumerate(methods):
        x = Inches(0.42) + i * (mw + gap)
        c = box(sl, x, my, mw, mh, CARD, RGBColor(0x2a,0x2a,0x2a), shadow=True)
        accent_bar(sl, x, my, mh, ORANGE)
        txt(sl, num, x+Inches(0.22), my+Pt(10), mw, Pt(18),
            size=9, color=DGRAY, font='Courier New')
        txt(sl, name, x+Inches(0.22), my+Pt(26), mw-Inches(0.3), Pt(54),
            size=17, color=WHITE, font='Georgia', wrap=True)
        txt(sl, desc, x+Inches(0.22), my+Pt(92), mw-Inches(0.3), Pt(74),
            size=11, color=GRAY, wrap=True)
        lv = sl.shapes.add_shape(1, x+Inches(0.22), my+Inches(2.28), Inches(1.02), Pt(21))
        fill_solid(lv, RGBColor(0x28,0x10,0x00)); line_solid(lv, RGBColor(0x80,0x38,0x00), 1)
        txt(sl, lvl, x+Inches(0.24), my+Inches(2.29), Inches(1.0), Pt(19),
            size=9, color=ORANGE, font='Courier New')
        pv = sl.shapes.add_shape(1, x+Inches(1.30), my+Inches(2.28), Inches(1.1), Pt(21))
        fill_solid(pv, CARD2); line_solid(pv, RGBColor(0x33,0x33,0x33), 1)
        txt(sl, pay, x+Inches(1.32), my+Inches(2.29), Inches(1.08), Pt(19),
            size=9, color=GRAY, font='Courier New')
        hbar(sl, x+Inches(0.22), my+Inches(2.57), mw-Inches(0.3), DGRAY)
        txt(sl, stack, x+Inches(0.22), my+Inches(2.63), mw-Inches(0.3), Pt(18),
            size=9, color=DGRAY, font='Courier New')

    txt(sl,
        'Start with 001–003. Run one method for 60 days before switching. Consistency beats cleverness.',
        Inches(0.42), Inches(5.25), Inches(12.4), Pt(28),
        size=13, color=GRAY, align=PP_ALIGN.CENTER)


def s05_toolkit(prs):
    sl = new_slide(prs)
    section_header(sl,
        '03 / Toolkit',
        'The working stack, ',
        'nothing more.',
        'Mastery over a small set beats shallow familiarity with twenty. Start with the first four.')

    tools = [
        ('Language model', 'ChatGPT',        'Daily driver for writing, research, ideation.'),
        ('Language model', 'Gemini',          'Long documents, multimodal input, analysis.'),
        ('Image gen',      'Midjourney',      'High-end visuals, brand & product art.'),
        ('Image gen',      'DALL·E',          'Quick illustrations, mockups, social content.'),
        ('Voice',          'ElevenLabs',      'Realistic narration in any language.'),
        ('Video',          'Pictory',         'Turn articles into short-form video.'),
        ('Design',         'Canva AI',        'Templates, social posts, branding.'),
        ('Workspace',      'Notion AI',       'Second brain — research, drafts, plans.'),
        ('Automation',     'Make',            'Visual workflows linking AI to your apps.'),
        ('Storefront',     'Gumroad',         'Sell digital products in under an hour.'),
        ('Marketplace',    'Etsy',            'Printables and AI art with built-in traffic.'),
        ('Freelance',      'Upwork / Fiverr', 'First five clients without an audience.'),
    ]

    tw, th = Inches(3.06), Inches(1.14)
    gx, gy = Inches(0.12), Inches(0.1)
    sx, sy = Inches(0.42), Inches(2.22)
    for i, (cat, name, purpose) in enumerate(tools):
        col, row = i % 4, i // 4
        x = sx + col * (tw + gx)
        y = sy + row * (th + gy)
        c = box(sl, x, y, tw, th, CARD, RGBColor(0x2a,0x2a,0x2a), shadow=True)
        txt(sl, cat,     x+Inches(0.18), y+Pt(10), tw-Inches(0.36), Pt(16),
            size=9, color=DGRAY, font='Courier New')
        txt(sl, name,    x+Inches(0.18), y+Pt(25), tw-Inches(0.36), Pt(34),
            size=20, color=WHITE, font='Georgia')
        txt(sl, purpose, x+Inches(0.18), y+Pt(64), tw-Inches(0.36), Pt(30),
            size=11, color=GRAY)


def s06_roadmap(prs):
    sl = new_slide(prs)
    section_header(sl,
        '04 / Roadmap',
        'A 60-day field plan, ',
        'step by step.',
        'Pick ONE method and run this protocol for two months before evaluating.')

    phases = [
        ('01', 'Week 1',   'Pick a single lane',
         'Match your skill to one method. Writing → 001/007. Visual → 002/003. Tech → 008/009.',
         ['List your top 3 skills', 'Pick one method. Close the other tabs.']),
        ('02', 'Wk 1–2',  'Master two tools',
         'Thirty minutes a day for 14 days. Build a personal prompt library as you go.',
         ['Pick your primary LLM', 'Save every prompt that worked']),
        ('03', 'Wk 2–3',  'Ship a first deliverable',
         'A prompt pack. A Notion template. Five sample articles. Proof you can deliver.',
         ['Define the smallest deliverable', 'Publish it somewhere public']),
        ('04', 'Wk 3–4',  'First three clients',
         'Trade price for testimonials. Three clients with written feedback unlocks real pricing.',
         ['Send 10 pitches per day', 'Collect a testimonial after every job']),
        ('05', 'Wk 5–6',  'Daily outreach + content',
         'One platform, one post format, daily. Ten outbound conversations a day.',
         ['Pick one platform: LinkedIn or X', 'Ten outbound conversations daily']),
        ('06', 'Wk 6–7',  'Sell outcomes, not hours',
         '"10 qualified leads/month" is worth more than "1,000-word article." Reframe every offer.',
         ['Rewrite offer as an outcome', 'Productize at three price points']),
        ('07', 'Wk 7–8',  'Layer in automation',
         'Once delivery is repeatable, wire it. Make, Zapier, prompt chains. Double income, same time.',
         ['Map your service as a flowchart', 'Automate the three slowest steps']),
    ]

    ph = Inches(0.715)
    for i, (num, week, title, desc, tasks) in enumerate(phases):
        y = Inches(1.9) + i * ph
        # phase number
        nb = sl.shapes.add_textbox(Inches(0.42), y, Inches(0.95), Inches(0.7))
        p = nb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = num
        r.font.name = 'Georgia'; r.font.size = Pt(50); r.font.color.rgb = LIME
        txt(sl, week,  Inches(0.42), y+Pt(54), Inches(1.1), Pt(16),
            size=9, color=DGRAY, font='Courier New')
        # title + desc + tasks
        txt(sl, title, Inches(1.62), y+Pt(4),  Inches(4.5), Pt(28),
            size=18, color=WHITE, font='Georgia', bold=False)
        txt(sl, desc,  Inches(1.62), y+Pt(28), Inches(4.6), Pt(30),
            size=11, color=GRAY, wrap=True)
        arrow_list(sl, tasks, Inches(6.5), y+Pt(4), Inches(6.4), Pt(58), size=11)
        hbar(sl, Inches(0.42), y+ph-Pt(1), Inches(12.4), RGBColor(0x22,0x22,0x22))


def s07_caveats(prs):
    sl = new_slide(prs)
    section_header(sl,
        '05 / Caveats',
        'What no one selling courses ',
        'will tell you.',
        'Real friction lives in places the influencers skip. Internalize these before spending a dollar.')

    warns = [
        ('Caveat 01', '"Easy money" is the red flag',
         'Anyone promising $10K/month passive without effort is selling a course about earning $10K/month passive. '
         'Real businesses take 3–6 months of focused work before meaningful revenue.'),
        ('Caveat 02', 'Fact-check every AI output',
         'Generative models confabulate. Numbers, quotes, and citations must be verified before reaching a client. '
         'The fastest way to destroy a reputation is one hallucinated statistic.'),
        ('Caveat 03', 'The market is flooded with slop',
         'Generic AI output is already a commodity. Winning means narrow niche + deep insight + production quality. '
         'Compete on taste and judgment, not throughput.'),
        ('Caveat 04', 'Read the licensing terms',
         'Commercial-use rights on AI-generated content shift frequently. '
         'Verify your right to sell before you list a single product.'),
    ]

    ww, wh = Inches(6.12), Inches(2.28)
    for i, (num, heading, body) in enumerate(warns):
        col, row = i % 2, i // 2
        x = Inches(0.42) + col * (ww + Inches(0.2))
        y = Inches(2.22) + row * (wh + Inches(0.14))
        c = box(sl, x+Pt(4), y, ww-Pt(4), wh, CARD, RGBColor(0x2a,0x2a,0x2a), shadow=True)
        bar = sl.shapes.add_shape(1, x, y, Pt(4), wh)
        fill_solid(bar, ORANGE); line_none(bar)
        txt(sl, num,     x+Pt(14), y+Pt(12), ww-Pt(20), Pt(18),
            size=9, color=ORANGE, font='Courier New', bold=True)
        txt(sl, heading, x+Pt(14), y+Pt(28), ww-Pt(20), Pt(38),
            size=20, color=WHITE, font='Georgia')
        txt(sl, body,    x+Pt(14), y+Pt(70), ww-Pt(20), Pt(82),
            size=12, color=GRAY, wrap=True)


def s08_closing(prs):
    sl = new_slide(prs)

    for i in range(14):
        s = sl.shapes.add_shape(1, Inches(i), Inches(0), Pt(1), H)
        fill_solid(s, RGBColor(0x1a,0x1a,0x1a)); line_none(s)

    tb = sl.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.8), Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'Start small. Ship daily. '
    r1.font.name = 'Georgia'; r1.font.size = Pt(52); r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = 'Compound monthly.'
    r2.font.name = 'Georgia'; r2.font.size = Pt(52)
    r2.font.color.rgb = LIME; r2.font.italic = True

    txt(sl,
        'This notebook is a living document — built to help think clearly about using AI '
        'to create real income. Updated as the landscape shifts.',
        Inches(0.75), Inches(4.6), Inches(9.5), Pt(46),
        size=15, color=GRAY, wrap=True)

    hbar(sl, Inches(0.75), Inches(5.72), Inches(11.8))

    secs = ['01 Market', '02 Methods', '03 Toolkit', '04 Roadmap', '05 Caveats']
    for i, s in enumerate(secs):
        txt(sl, s, Inches(0.75)+i*Inches(2.45), Inches(5.84),
            Inches(2.3), Pt(20), size=10, color=DGRAY, font='Courier New')

    txt(sl, '// AI Income Lab — Personal study site  ·  Built May 2026  ·  v0.2',
        Inches(0.75), Inches(7.08), Inches(11.8), Pt(20),
        size=9, color=DGRAY, font='Courier New', align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    steps = [
        ('Title',          s01_title),
        ('Market / Stats', s02_market),
        ('Methods 1–5',    s03_methods_a),
        ('Methods 6–10',   s04_methods_b),
        ('Toolkit',        s05_toolkit),
        ('Roadmap',        s06_roadmap),
        ('Caveats',        s07_caveats),
        ('Closing',        s08_closing),
    ]
    for i, (name, fn) in enumerate(steps, 1):
        print(f'  [{i}/{len(steps)}] {name}...')
        fn(prs)

    out = '/home/user/yupebis/AI_Income_Lab_v2.pptx'
    prs.save(out)
    print(f'\n✅  Saved → {out}')

if __name__ == '__main__':
    main()
