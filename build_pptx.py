"""
Avengers: Doomsday — 3D Copper Glow Presentation builder
Generates a 15-slide .pptx with dark stage / copper neon aesthetic.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy
import math

# ── Colour palette ──────────────────────────────────────────────────────────
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_PANEL  = RGBColor(0x11, 0x11, 0x11)
COPPER      = RGBColor(0xC8, 0x79, 0x41)
AMBER       = RGBColor(0xE8, 0x92, 0x3A)
ORANGE_GLOW = RGBColor(0xFF, 0x78, 0x20)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CREAM       = RGBColor(0xF5, 0xE6, 0xD0)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
GOLD        = RGBColor(0xFF, 0xD7, 0x00)
PURPLE      = RGBColor(0x7B, 0x1F, 0xA2)
RED_GLOW    = RGBColor(0xCC, 0x00, 0x00)
DARK_RED    = RGBColor(0x1A, 0x00, 0x00)

# ── Slide dimensions ─────────────────────────────────────────────────────────
W = Inches(13.3)
H = Inches(7.5)

# ── Helper: EMU shortcuts ────────────────────────────────────────────────────
def inch(v): return Inches(v)

# ── Low-level XML helpers ────────────────────────────────────────────────────

def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def set_shape_fill(shape, color: RGBColor, alpha: int = 255):
    """Solid fill with optional transparency (alpha 0-255)."""
    sp = shape.shape._element if hasattr(shape, 'shape') else shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    # remove existing fill
    for tag in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'), qn('a:pattFill')]:
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', rgb_hex(color))
    if alpha < 255:
        lumMod = etree.SubElement(srgbClr, qn('a:alpha'))
        lumMod.set('val', str(int(alpha / 255 * 100000)))


def set_no_fill(shape):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    for tag_no in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill')]:
        el = spPr.find(tag_no)
        if el is not None:
            spPr.remove(el)
    etree.SubElement(spPr, qn('a:noFill'))


def set_line(shape, color: RGBColor, width_pt: float = 2, alpha: int = 255):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(int(width_pt * 12700)))
    solidFill = etree.SubElement(ln, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', rgb_hex(color))
    if alpha < 255:
        a = etree.SubElement(srgbClr, qn('a:alpha'))
        a.set('val', str(int(alpha / 255 * 100000)))


def set_no_line(shape):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:noFill'))


def add_shadow(shape, blur_pt=28, offset_pt=5, angle=135, opacity=0.75):
    """Outer shadow effect."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(int(blur_pt * 12700)))
    outerShdw.set('dist', str(int(offset_pt * 12700)))
    outerShdw.set('dir', str(int(angle * 60000)))
    outerShdw.set('algn', 'ctr')
    outerShdw.set('rotWithShape', '0')
    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', '000000')
    alphaEl = etree.SubElement(srgbClr, qn('a:alpha'))
    alphaEl.set('val', str(int(opacity * 100000)))


def set_grad_fill(shape, stops):
    """
    stops = list of (pos_pct, RGBColor, alpha_pct)
    Linear gradient left→right (or top→bottom with ang=5400000).
    """
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    for tag_gf in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill')]:
        el = spPr.find(tag_gf)
        if el is not None:
            spPr.remove(el)
    gradFill = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    for pos, color, alpha_pct in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        srgbClr = etree.SubElement(gs, qn('a:srgbClr'))
        srgbClr.set('val', rgb_hex(color))
        if alpha_pct < 100:
            a = etree.SubElement(srgbClr, qn('a:alpha'))
            a.set('val', str(int(alpha_pct * 1000)))
    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', '5400000')   # top→bottom
    lin.set('scaled', '0')


def set_text_glow(run, color: RGBColor, radius_pt=8, alpha=0.7):
    """Glow effect on a text run element."""
    rPr = run._r
    effectLst = rPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = etree.SubElement(rPr, qn('a:effectLst'))
    glow = etree.SubElement(effectLst, qn('a:glow'))
    glow.set('rad', str(int(radius_pt * 12700)))
    srgbClr = etree.SubElement(glow, qn('a:srgbClr'))
    srgbClr.set('val', rgb_hex(color))
    a = etree.SubElement(srgbClr, qn('a:alpha'))
    a.set('val', str(int(alpha * 100000)))


def set_shape_glow(shape, color: RGBColor, radius_pt=20, alpha=0.65):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    glow = etree.SubElement(effectLst, qn('a:glow'))
    glow.set('rad', str(int(radius_pt * 12700)))
    srgbClr = etree.SubElement(glow, qn('a:srgbClr'))
    srgbClr.set('val', rgb_hex(color))
    a = etree.SubElement(srgbClr, qn('a:alpha'))
    a.set('val', str(int(alpha * 100000)))


def set_rotation(shape, degrees: float):
    """Rotate a shape by given degrees (positive = clockwise)."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    xfrm = spPr.find(qn('a:xfrm')) if spPr is not None else None
    if xfrm is None and spPr is not None:
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    if xfrm is not None:
        xfrm.set('rot', str(int(degrees * 60000)))


# ── Slide factory helpers ────────────────────────────────────────────────────

def new_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    slide = prs.slides.add_slide(blank_layout)
    return slide


def black_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def add_floor_glow(slide, accent=COPPER):
    """Horizontal copper glow strip at ~60% height simulating reflective floor."""
    # gradient rect spanning full width, bottom 40%
    rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        inch(0), inch(4.5),
        W, inch(3.0)
    )
    set_grad_fill(rect, [
        (0,   accent,                5),
        (15,  RGBColor(0x30,0x18,0x00), 40),
        (100, BLACK,                100),
    ])
    set_no_line(rect)

    # thin bright copper line
    line = slide.shapes.add_shape(1, inch(0), inch(4.52), W, Pt(2))
    set_shape_fill(line, accent, 220)
    set_no_line(line)
    set_shape_glow(line, accent, radius_pt=12, alpha=0.8)


def add_light_beams(slide, color=COPPER):
    """Two diagonal volumetric light shafts."""
    for x, y, w, h, rot in [
        (inch(8),  inch(-1), inch(3), inch(10), 30),
        (inch(2),  inch(-1), inch(2), inch(9),  -28),
    ]:
        beam = slide.shapes.add_shape(1, x, y, w, h)
        set_shape_fill(beam, color, int(0.12 * 255))
        set_no_line(beam)
        set_rotation(beam, rot)


def add_glow_halo(slide, cx, cy, rx, ry, color=ORANGE_GLOW, alpha=80):
    """Soft radial ellipse for backlit neon effect."""
    oval = slide.shapes.add_shape(
        9,   # MSO_SHAPE_TYPE.OVAL
        cx - rx, cy - ry, rx*2, ry*2
    )
    set_shape_fill(oval, color, alpha)
    set_no_line(oval)


def add_depth_panels(slide, accent=COPPER):
    """Two blurred background depth panels offset left and right."""
    for x, y, w, h in [
        (inch(0.3), inch(1.0), inch(4.0), inch(5.0)),
        (inch(9.2), inch(1.2), inch(3.8), inch(4.8)),
    ]:
        panel = slide.shapes.add_shape(1, x, y, w, h)
        set_shape_fill(panel, DARK_PANEL, int(0.55 * 255))
        set_line(panel, accent, 1, int(0.5 * 255))
        add_shadow(panel, blur_pt=28, offset_pt=5, opacity=0.75)


def add_main_card(slide, x, y, w, h, accent=COPPER):
    """Main floating glass card — dark with copper border & shadow."""
    card = slide.shapes.add_shape(1, x, y, w, h)
    set_shape_fill(card, DARK_PANEL, int(0.85 * 255))
    set_line(card, accent, 2)
    add_shadow(card, blur_pt=28, offset_pt=5, opacity=0.75)
    return card


def add_separator_line(slide, x, y, w, color=COPPER):
    line = slide.shapes.add_shape(1, x, y, w, Pt(1))
    set_shape_fill(line, color, 220)
    set_no_line(line)


def add_title_text(slide, text, x, y, w, h, size=52, color=WHITE, bold=True,
                   font='Impact', align=PP_ALIGN.CENTER, glow=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    if glow:
        set_text_glow(run, AMBER, radius_pt=10, alpha=0.85)
    return txBox


def add_body_text(slide, text, x, y, w, h, size=14, color=LIGHT_GRAY,
                  font='Calibri', align=PP_ALIGN.LEFT, bold=False, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_bullet_list(slide, items, x, y, w, h, size=14, bullet_char='◆ '):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.name = 'Calibri'
        run.font.size = Pt(size)
        run.font.color.rgb = LIGHT_GRAY


def add_copper_bullet(slide, items, x, y, w, h, size=14):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        # bullet icon in copper
        r1 = p.add_run()
        r1.text = '◆ '
        r1.font.name = 'Calibri'
        r1.font.size = Pt(size)
        r1.font.color.rgb = AMBER
        r1.font.bold = True
        # text in light gray
        r2 = p.add_run()
        r2.text = item
        r2.font.name = 'Calibri'
        r2.font.size = Pt(size)
        r2.font.color.rgb = LIGHT_GRAY


def add_stat_card(slide, x, y, w, h, number, label,
                  accent=AMBER, num_size=54, lbl_size=13):
    card = add_main_card(slide, x, y, w, h, accent)
    # big stat number
    nb = slide.shapes.add_textbox(x, y + inch(0.15), w, inch(0.9))
    p = nb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number
    r.font.name = 'Impact'
    r.font.size = Pt(num_size)
    r.font.color.rgb = accent
    set_text_glow(r, accent, radius_pt=8, alpha=0.8)
    # label
    lb = slide.shapes.add_textbox(x, y + inch(0.95), w, inch(0.55))
    p2 = lb.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = 'Calibri'
    r2.font.size = Pt(lbl_size)
    r2.font.color.rgb = WHITE


def add_char_card(slide, x, y, w, h, char_name, actor_name, accent=AMBER):
    card = add_main_card(slide, x, y, w, h, accent)
    # character name — white bold
    nb = slide.shapes.add_textbox(x + inch(0.1), y + h - inch(0.8), w - inch(0.2), inch(0.42))
    p = nb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = char_name
    r.font.name = 'Calibri'
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    r.font.bold = True
    # actor name — copper italic
    ab = slide.shapes.add_textbox(x + inch(0.1), y + h - inch(0.42), w - inch(0.2), inch(0.4))
    p2 = ab.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = actor_name
    r2.font.name = 'Trebuchet MS'
    r2.font.size = Pt(11)
    r2.font.color.rgb = AMBER
    r2.font.italic = True
    return card


def add_info_card(slide, x, y, w, h, text, accent=AMBER):
    card = add_main_card(slide, x, y, w, h, accent)
    tb = slide.shapes.add_textbox(x + inch(0.1), y + inch(0.12), w - inch(0.2), h - inch(0.25))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = 'Calibri'
    r.font.size = Pt(13)
    r.font.color.rgb = LIGHT_GRAY
    return card


def slide_base(slide, accent=COPPER):
    """Apply black bg, floor glow, light beams, and depth panels."""
    black_bg(slide)
    add_floor_glow(slide, accent)
    add_light_beams(slide, accent)
    add_depth_panels(slide, accent)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════

def build_slide_1(prs):
    """TITLE SLIDE"""
    slide = new_slide(prs)
    slide_base(slide)

    # Glow halo behind logo
    add_glow_halo(slide, inch(6.65), inch(2.4), inch(2.2), inch(1.4), ORANGE_GLOW, 90)

    # Pedestal ring glow
    add_glow_halo(slide, inch(6.65), inch(4.2), inch(2.8), inch(0.55), COPPER, 130)

    # "A" logo placeholder — stylised big A in copper
    logo = slide.shapes.add_textbox(inch(5.4), inch(1.0), inch(2.5), inch(2.2))
    p = logo.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Ⓐ'
    r.font.name = 'Impact'
    r.font.size = Pt(96)
    r.font.color.rgb = AMBER
    set_text_glow(r, ORANGE_GLOW, radius_pt=22, alpha=0.9)

    # Main title
    add_title_text(slide, 'AVENGERS: DOOMSDAY',
                   inch(1.5), inch(3.5), inch(10.3), inch(1.3),
                   size=56, glow=True)

    # Subtitle
    sub = slide.shapes.add_textbox(inch(2.0), inch(4.7), inch(9.3), inch(0.7))
    p = sub.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"The Beginning of the End of the Multiverse Saga"'
    r.font.name = 'Trebuchet MS'
    r.font.size = Pt(20)
    r.font.color.rgb = CREAM
    r.font.italic = True

    # Bottom line
    add_body_text(slide, 'Marvel Studios  ·  Phase 6  ·  December 18, 2026',
                  inch(3.0), inch(6.5), inch(7.3), inch(0.5),
                  size=14, color=COPPER, align=PP_ALIGN.CENTER)


def build_slide_2(prs):
    """WHAT IS THIS FILM?"""
    slide = new_slide(prs)
    slide_base(slide)

    # Main card center-left
    card = add_main_card(slide, inch(0.6), inch(0.8), inch(7.8), inch(5.5))

    add_title_text(slide, 'WHAT IS THIS FILM?',
                   inch(0.9), inch(0.9), inch(7.2), inch(0.9), size=36)
    add_separator_line(slide, inch(0.9), inch(1.85), inch(7.2))

    # Bullet points
    bullets = [
        'The 39th film in the Marvel Cinematic Universe',
        'Phase 6 — the penultimate chapter of the Multiverse Saga',
        'Robert Downey Jr. returns — not as Tony Stark, but as Doctor Doom',
        'Directed by the Russo Brothers — returning after Avengers: Endgame',
    ]
    add_copper_bullet(slide, bullets, inch(1.0), inch(2.1), inch(7.0), inch(3.5), size=16)

    # Stat bar
    bar = add_main_card(slide, inch(0.6), inch(6.3), inch(12.0), inch(0.7), COPPER)
    add_body_text(slide, '39th MCU Film   ·   Phase 6   ·   Sequel: Secret Wars 2027',
                  inch(0.8), inch(6.35), inch(11.6), inch(0.55),
                  size=15, color=AMBER, align=PP_ALIGN.CENTER, bold=True)


def build_slide_3(prs):
    """THE VILLAIN: DOCTOR DOOM"""
    slide = new_slide(prs)
    slide_base(slide, RED_GLOW)  # red variant

    add_glow_halo(slide, inch(6.65), inch(3.5), inch(2.5), inch(1.6), RED_GLOW, 110)

    card = add_main_card(slide, inch(1.5), inch(0.7), inch(9.8), inch(5.6), RED_GLOW)

    add_title_text(slide, 'THE VILLAIN: DOCTOR DOOM',
                   inch(1.8), inch(0.8), inch(9.2), inch(0.9), size=36, glow=False)
    t = slide.shapes[-1].text_frame.paragraphs[0].runs[0]
    set_text_glow(t, RED_GLOW, radius_pt=10, alpha=0.85)

    add_title_text(slide, 'VICTOR VON DOOM',
                   inch(1.8), inch(1.75), inch(9.2), inch(0.8), size=44, color=AMBER)

    add_separator_line(slide, inch(1.8), inch(2.6), inch(9.2), RED_GLOW)

    add_body_text(slide,
        'Doctor Doom — ruler of Latveria and master of both science and sorcery — '
        'emerges as the MCU\'s most formidable villain yet. His command over the '
        'multiverse threatens every reality simultaneously.',
        inch(1.9), inch(2.7), inch(9.0), inch(1.8), size=15)

    # Three info cards
    for i, (title, desc) in enumerate([
        ("RDJ's Return", "Robert Downey Jr. returns in a shocking new role"),
        ("Reality-Bending Threat", "Doom controls the fabric of the multiverse"),
        ("Replaced Kang", "Doom replaces Kang as the saga's prime villain"),
    ]):
        cx = inch(1.8) + i * inch(3.3)
        add_info_card(slide, cx, inch(4.7), inch(3.1), inch(1.35), f"{title}\n{desc}", RED_GLOW)


def build_slide_4(prs):
    """RDJ: FROM IRON MAN TO IRON DOOM"""
    slide = new_slide(prs)
    slide_base(slide)

    # Left panel — RDJ character info
    left = add_main_card(slide, inch(0.5), inch(0.8), inch(5.5), inch(5.6))

    add_title_text(slide, 'RDJ: FROM IRON MAN\nTO IRON DOOM',
                   inch(0.7), inch(0.9), inch(5.0), inch(1.3), size=30)
    add_separator_line(slide, inch(0.7), inch(2.3), inch(5.0))

    add_body_text(slide,
        'Robert Downey Jr. portrayed Tony Stark / Iron Man across 10 films '
        'from 2008 to 2019. His sacrifice in Avengers: Endgame defined a generation '
        'of cinema.\n\n'
        'Now, in a multiverse-shattering twist, RDJ returns — not as the hero '
        'we knew, but as Victor Von Doom: the most powerful villain in Marvel history.',
        inch(0.7), inch(2.4), inch(5.0), inch(3.0), size=14)

    # Right side
    right = add_main_card(slide, inch(6.4), inch(0.8), inch(6.3), inch(5.6))

    # Pull quote
    quote = slide.shapes.add_textbox(inch(6.6), inch(1.1), inch(5.9), inch(2.5))
    p = quote.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"He saved us once.\nNow he threatens\neverything."'
    r.font.name = 'Impact'
    r.font.size = Pt(30)
    r.font.color.rgb = AMBER
    r.font.italic = True
    set_text_glow(r, AMBER, radius_pt=8, alpha=0.8)

    # Timeline nodes
    nodes = [
        ('Iron Man', '2008'),
        ('Endgame', '2019'),
        ('Doctor Doom', '2026'),
    ]
    timeline_y = inch(5.4)
    for i, (label, year) in enumerate(nodes):
        nx = inch(6.9) + i * inch(1.9)
        dot = slide.shapes.add_shape(9, nx - inch(0.2), timeline_y - inch(0.2),
                                      inch(0.4), inch(0.4))
        set_shape_fill(dot, AMBER, 255)
        set_no_line(dot)
        set_shape_glow(dot, AMBER, radius_pt=8, alpha=0.9)
        if i < 2:
            ln = slide.shapes.add_shape(1, nx + inch(0.2), timeline_y - inch(0.04),
                                          inch(1.5), inch(0.08))
            set_shape_fill(ln, COPPER, 200)
            set_no_line(ln)
        add_body_text(slide, year, nx - inch(0.5), timeline_y + inch(0.25),
                      inch(1.0), inch(0.3), size=11, color=AMBER, align=PP_ALIGN.CENTER)
        add_body_text(slide, label, nx - inch(0.6), timeline_y + inch(0.55),
                      inch(1.2), inch(0.3), size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def build_slide_5(prs):
    """THE RUSSO BROTHERS"""
    slide = new_slide(prs)
    slide_base(slide)

    # Two floating panels side by side
    left = add_main_card(slide, inch(0.5), inch(0.8), inch(6.0), inch(5.8))
    right = add_main_card(slide, inch(6.9), inch(0.8), inch(5.9), inch(5.8))

    add_title_text(slide, 'THE RUSSO BROTHERS',
                   inch(0.7), inch(0.9), inch(5.6), inch(0.85), size=28)
    add_separator_line(slide, inch(0.7), inch(1.82), inch(5.6))

    add_body_text(slide,
        'Anthony and Joe Russo are the directors behind the two highest-grossing '
        'superhero films of all time. After a 7-year hiatus from the MCU, they '
        'return to close the Multiverse Saga.',
        inch(0.7), inch(2.0), inch(5.6), inch(2.5), size=14)

    # Copper divider
    add_separator_line(slide, inch(6.7), inch(1.0), inch(0.05), COPPER)

    # Stat cards on the right
    stats = [
        ('4', 'MCU Directorial Credits'),
        ('$2.79B', 'Endgame Worldwide Gross'),
        ('#1', 'Highest-Grossing Superhero Film Ever'),
        ('2027', 'Also Directing Secret Wars'),
    ]
    for i, (num, lbl) in enumerate(stats):
        sy = inch(1.0) + i * inch(1.28)
        add_stat_card(slide, inch(7.0), sy, inch(5.6), inch(1.1),
                      num, lbl, num_size=36, lbl_size=12)


def build_slide_6(prs):
    """THE HEROES — 3×3 grid"""
    slide = new_slide(prs)
    slide_base(slide)

    add_title_text(slide, 'THE HEROES',
                   inch(1.0), inch(0.2), inch(11.3), inch(0.75), size=44)

    characters = [
        ('Captain America', 'Anthony Mackie'),
        ('Thor', 'Chris Hemsworth'),
        ('Spider-Man', 'Tom Holland'),
        ('Shuri', 'Letitia Wright'),
        ('Loki', 'Tom Hiddleston'),
        ('Ant-Man', 'Paul Rudd'),
        ('Yelena Belova', 'Florence Pugh'),
        ('Sentry', 'Lewis Pullman'),
        ('Namor', 'Tenoch Huerta'),
    ]

    cols, rows = 3, 3
    cw, ch = inch(4.0), inch(1.95)
    sx, sy = inch(0.55), inch(1.1)
    gap_x, gap_y = inch(0.15), inch(0.12)

    for i, (char, actor) in enumerate(characters):
        col = i % cols
        row = i // cols
        x = sx + col * (cw + gap_x)
        y = sy + row * (ch + gap_y)
        add_char_card(slide, x, y, cw, ch, char, actor)


def build_slide_7(prs):
    """THE FANTASTIC FOUR — gold variant"""
    slide = new_slide(prs)
    slide_base(slide, GOLD)

    add_glow_halo(slide, inch(6.65), inch(3.0), inch(3.5), inch(1.8), GOLD, 85)

    add_title_text(slide, 'THE FANTASTIC FOUR',
                   inch(1.0), inch(0.2), inch(11.3), inch(0.85), size=46, color=GOLD)

    # Large central promo card
    promo = add_main_card(slide, inch(1.5), inch(1.1), inch(9.7), inch(3.5), GOLD)
    add_body_text(slide,
        'Reed Richards · Sue Storm · Human Torch · The Thing\n\n'
        'The First Family of Marvel finally joins the MCU, arriving at the '
        'critical moment when Doctor Doom — once Reed\'s closest rival — '
        'threatens all of existence.',
        inch(1.8), inch(1.4), inch(9.1), inch(2.9),
        size=16, color=CREAM, align=PP_ALIGN.CENTER)

    # Four character cards
    ff = [
        ('Reed Richards', 'Pedro Pascal'),
        ('Sue Storm', 'Vanessa Kirby'),
        ('Human Torch', 'Joseph Quinn'),
        ('The Thing', 'Ebon Moss-Bachrach'),
    ]
    cw, ch = inch(3.05), inch(1.5)
    for i, (char, actor) in enumerate(ff):
        x = inch(0.35) + i * (cw + inch(0.12))
        add_char_card(slide, x, inch(4.85), cw, ch, char, actor, GOLD)


def build_slide_8(prs):
    """THE X-MEN ARRIVE — purple variant"""
    slide = new_slide(prs)
    slide_base(slide, PURPLE)

    add_glow_halo(slide, inch(6.65), inch(3.2), inch(2.8), inch(1.6), PURPLE, 100)

    add_title_text(slide, 'THE X-MEN ARRIVE',
                   inch(1.0), inch(0.2), inch(11.3), inch(0.85), size=46, color=PURPLE)

    # Central icon card
    icon_card = add_main_card(slide, inch(4.5), inch(1.1), inch(4.3), inch(2.2), PURPLE)
    icon_txt = slide.shapes.add_textbox(inch(4.7), inch(1.3), inch(3.9), inch(1.8))
    p = icon_txt.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Ⓧ'
    r.font.name = 'Impact'
    r.font.size = Pt(72)
    r.font.color.rgb = PURPLE
    set_text_glow(r, PURPLE, radius_pt=18, alpha=0.9)

    # Two character cards
    chars = [
        ('Professor X', 'Patrick Stewart'),
        ('Beast', 'Kelsey Grammer'),
    ]
    for i, (char, actor) in enumerate(chars):
        x = inch(1.2) + i * inch(6.5)
        add_char_card(slide, x, inch(3.6), inch(4.5), inch(1.9), char, actor, PURPLE)

    # Pull quote
    q = slide.shapes.add_textbox(inch(1.5), inch(5.65), inch(10.3), inch(0.9))
    p = q.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"The Multiverse makes everything possible."'
    r.font.name = 'Trebuchet MS'
    r.font.size = Pt(22)
    r.font.color.rgb = PURPLE
    r.font.italic = True
    set_text_glow(r, PURPLE, radius_pt=8, alpha=0.8)


def build_slide_9(prs):
    """THE PLOT"""
    slide = new_slide(prs)
    slide_base(slide)

    card = add_main_card(slide, inch(0.6), inch(0.7), inch(12.1), inch(5.7))

    add_title_text(slide, 'THE PLOT',
                   inch(0.9), inch(0.8), inch(11.5), inch(0.85), size=40)
    add_separator_line(slide, inch(0.9), inch(1.75), inch(11.5))

    add_body_text(slide,
        'When the boundaries between parallel universes begin to collapse, '
        'Victor Von Doom — the brilliant yet ruthless ruler of Latveria — '
        'seizes control of the fractures, weaponising the multiverse itself.\n\n'
        'Earth\'s Mightiest Heroes, aided by the newly arrived Fantastic Four '
        'and long-lost mutants, must unite across realities to prevent Doom '
        'from rewriting existence in his own image.\n\n'
        'Spanning three universes and bringing together characters from the '
        'entire 18-year MCU saga, Avengers: Doomsday is the largest convergence '
        'event in cinema history — a clash that will permanently reshape the '
        'Marvel Multiverse and set the stage for Secret Wars in 2027.',
        inch(1.0), inch(1.9), inch(11.1), inch(3.3), size=16, color=CREAM)

    # Three stat callouts
    for i, (num, lbl) in enumerate([
        ('3', 'Universes'), ('39', 'Films Leading Here'), ('1', 'Villain to Rule Them All')
    ]):
        cx = inch(1.2) + i * inch(3.9)
        add_stat_card(slide, cx, inch(5.9), inch(3.5), inch(1.15),
                      num, lbl, num_size=40, lbl_size=13)


def build_slide_10(prs):
    """THE THUNDERBOLTS"""
    slide = new_slide(prs)
    slide_base(slide)

    add_title_text(slide, 'THE THUNDERBOLTS',
                   inch(1.0), inch(0.2), inch(11.3), inch(0.82), size=44)

    sub = slide.shapes.add_textbox(inch(2.0), inch(1.02), inch(9.3), inch(0.5))
    p = sub.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"The Government\'s Secret Weapon"'
    r.font.name = 'Trebuchet MS'
    r.font.size = Pt(18)
    r.font.color.rgb = CREAM
    r.font.italic = True

    chars = [
        ('Yelena Belova', 'Florence Pugh'),
        ('Red Guardian', 'David Harbour'),
        ('U.S. Agent', 'Wyatt Russell'),
        ('Ghost', 'Hannah John-Kamen'),
        ('Sentry', 'Lewis Pullman'),
        ('Winter Soldier', 'Sebastian Stan'),
    ]

    cw, ch = inch(3.95), inch(2.35)
    for i, (char, actor) in enumerate(chars):
        col = i % 3
        row = i // 3
        x = inch(0.5) + col * (cw + inch(0.15))
        y = inch(1.65) + row * (ch + inch(0.15))
        add_char_card(slide, x, y, cw, ch, char, actor)


def build_slide_11(prs):
    """RELEASE DATE"""
    slide = new_slide(prs)
    slide_base(slide)

    add_title_text(slide, 'RELEASE DATE',
                   inch(1.0), inch(0.2), inch(11.3), inch(0.82), size=44)

    # Left panel — Avengers Doomsday
    lp = add_main_card(slide, inch(0.4), inch(1.1), inch(5.7), inch(4.6))
    add_title_text(slide, 'AVENGERS:\nDOOMSDAY',
                   inch(0.6), inch(1.3), inch(5.3), inch(1.5), size=30, glow=True)
    add_body_text(slide, 'Marvel Studios', inch(0.6), inch(2.9), inch(5.3), inch(0.4),
                  size=14, color=COPPER, align=PP_ALIGN.CENTER)
    add_body_text(slide, 'Phase 6 · Multiverse Saga Finale', inch(0.6), inch(3.3),
                  inch(5.3), inch(0.4), size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Centre date callout
    add_glow_halo(slide, inch(6.65), inch(3.0), inch(1.5), inch(1.0), ORANGE_GLOW, 80)
    date_box = slide.shapes.add_textbox(inch(5.0), inch(2.2), inch(3.3), inch(1.6))
    p = date_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'DECEMBER\n18, 2026'
    r.font.name = 'Impact'
    r.font.size = Pt(32)
    r.font.color.rgb = AMBER
    set_text_glow(r, AMBER, radius_pt=12, alpha=0.9)

    # Right panel — Dune Part Three
    rp = add_main_card(slide, inch(7.2), inch(1.1), inch(5.7), inch(4.6),
                       RGBColor(0x70, 0x70, 0x70))
    add_title_text(slide, 'DUNE:\nPART THREE',
                   inch(7.4), inch(1.3), inch(5.3), inch(1.5), size=30,
                   color=RGBColor(0xAA, 0xAA, 0xAA), glow=False)
    add_body_text(slide, 'Warner Bros. Pictures', inch(7.4), inch(2.9), inch(5.3), inch(0.4),
                  size=14, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)
    add_body_text(slide, 'Same Opening Weekend', inch(7.4), inch(3.3),
                  inch(5.3), inch(0.4), size=13, color=RGBColor(0x77,0x77,0x77), align=PP_ALIGN.CENTER)

    # Warning strip
    warn = add_main_card(slide, inch(0.4), inch(6.1), inch(12.5), inch(0.75), AMBER)
    add_body_text(slide, '⚠  DunesDay 2026 — Battle for IMAX Screens',
                  inch(0.6), inch(6.18), inch(12.1), inch(0.55),
                  size=16, color=AMBER, align=PP_ALIGN.CENTER, bold=True)


def build_slide_12(prs):
    """PHASE 6 TIMELINE"""
    slide = new_slide(prs)
    slide_base(slide)

    add_title_text(slide, 'PHASE 6 TIMELINE',
                   inch(1.0), inch(0.2), inch(11.3), inch(0.82), size=44)

    # Timeline line
    tl = slide.shapes.add_shape(1, inch(0.7), inch(3.55), inch(11.9), inch(0.12))
    set_shape_fill(tl, COPPER, 230)
    set_no_line(tl)
    set_shape_glow(tl, COPPER, radius_pt=6, alpha=0.8)

    nodes = [
        ('Brave New\nWorld', 'Feb 2025'),
        ('Thunderbolts*', 'May 2025'),
        ('Fantastic\nFour', 'Jul 2025'),
        ('Spider-Man:\nBrand New Day', 'Jul 2026'),
        ('Avengers:\nDoomsday', 'Dec 2026'),
        ('Secret\nWars', '2027'),
    ]

    for i, (label, date) in enumerate(nodes):
        nx = inch(0.9) + i * inch(2.35)
        ny = inch(3.5)
        is_main = i == 4
        r = inch(0.28) if is_main else inch(0.2)
        dot = slide.shapes.add_shape(9, nx - r, ny - r, r*2, r*2)
        col = AMBER if is_main else COPPER
        set_shape_fill(dot, col, 255)
        set_no_line(dot)
        set_shape_glow(dot, col, radius_pt=10 if is_main else 6, alpha=0.85)

        # Label above
        lbl = slide.shapes.add_textbox(nx - inch(1.1), ny - inch(2.0), inch(2.2), inch(1.8))
        p = lbl.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r2 = p.add_run()
        r2.text = ('⭐ ' if is_main else '') + label
        r2.font.name = 'Calibri'
        r2.font.size = Pt(13 if is_main else 12)
        r2.font.color.rgb = AMBER if is_main else LIGHT_GRAY
        r2.font.bold = is_main

        # Date below
        dt = slide.shapes.add_textbox(nx - inch(1.0), ny + inch(0.45), inch(2.0), inch(0.45))
        p2 = dt.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r3 = p2.add_run()
        r3.text = date
        r3.font.name = 'Calibri'
        r3.font.size = Pt(11)
        r3.font.color.rgb = COPPER


def build_slide_13(prs):
    """BY THE NUMBERS — 3×2 grid of stat cards"""
    slide = new_slide(prs)
    slide_base(slide)

    add_title_text(slide, 'BY THE NUMBERS',
                   inch(1.0), inch(0.2), inch(11.3), inch(0.82), size=44)

    stats = [
        ('39',     'MCU Film Number'),
        ('$2.79B', 'Endgame Benchmark'),
        ('4',      'Russo Brothers MCU Films'),
        ('3',      'Universes Colliding'),
        ('2026',   'Release Year'),
        ('1',      'Victor Von Doom'),
    ]

    cw, ch = inch(4.1), inch(1.75)
    for i, (num, lbl) in enumerate(stats):
        col = i % 3
        row = i // 3
        x = inch(0.4) + col * (cw + inch(0.15))
        y = inch(1.15) + row * (ch + inch(0.18))
        add_stat_card(slide, x, y, cw, ch, num, lbl, num_size=52, lbl_size=14)


def build_slide_14(prs):
    """WHY IT MATTERS"""
    slide = new_slide(prs)
    slide_base(slide)

    card = add_main_card(slide, inch(0.6), inch(0.7), inch(12.1), inch(4.5))

    add_title_text(slide, 'WHY IT MATTERS',
                   inch(0.9), inch(0.8), inch(11.5), inch(0.85), size=40)
    add_separator_line(slide, inch(0.9), inch(1.75), inch(11.5))

    add_body_text(slide,
        'Avengers: Doomsday is more than a superhero film — it is the culmination '
        'of an 18-year cinematic experiment that has redefined Hollywood blockbusters.\n\n'
        'Robert Downey Jr\'s transformation from beloved hero to ultimate villain '
        'represents one of the boldest casting decisions in franchise history. '
        'The return of the Russo Brothers guarantees the large-scale, emotionally '
        'resonant storytelling that defined the Infinity Saga.\n\n'
        'With the Fantastic Four, X-Men, Thunderbolts, and over 30 returning heroes, '
        'this film doesn\'t just close a chapter — it opens a new era.',
        inch(1.0), inch(1.9), inch(11.1), inch(2.9), size=15, color=CREAM)

    # Three icon cards at bottom
    icon_items = [
        ('🦸', 'MCU Veterans\n+ New Heroes'),
        ('🌌', 'Multiverse Saga\nEndgame'),
        ('🎭', "RDJ's Most\nComplex Role"),
    ]
    cw, ch = inch(3.8), inch(1.55)
    for i, (icon, lbl) in enumerate(icon_items):
        x = inch(0.7) + i * (cw + inch(0.35))
        ic = add_main_card(slide, x, inch(5.45), cw, ch)
        tb = slide.shapes.add_textbox(x + inch(0.1), inch(5.5), cw - inch(0.2), ch - inch(0.1))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = icon + '  ' + lbl
        r.font.name = 'Calibri'
        r.font.size = Pt(14)
        r.font.color.rgb = LIGHT_GRAY


def build_slide_15(prs):
    """CLOSING — DOOM IS COMING"""
    slide = new_slide(prs)
    black_bg(slide)
    add_floor_glow(slide)
    add_light_beams(slide)

    # Massive glow halo
    add_glow_halo(slide, inch(6.65), inch(4.1), inch(4.0), inch(1.2), ORANGE_GLOW, 100)
    add_glow_halo(slide, inch(6.65), inch(3.2), inch(3.0), inch(2.0), COPPER, 60)

    # Main title — DOOM IS COMING
    doom = slide.shapes.add_textbox(inch(0.5), inch(2.0), inch(12.3), inch(2.5))
    p = doom.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'DOOM IS COMING.'
    r.font.name = 'Impact'
    r.font.size = Pt(80)
    r.font.color.rgb = AMBER
    set_text_glow(r, ORANGE_GLOW, radius_pt=28, alpha=0.95)

    add_body_text(slide, 'December 18, 2026  ·  Marvel Studios',
                  inch(2.0), inch(4.8), inch(9.3), inch(0.6),
                  size=18, color=COPPER, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    builders = [
        build_slide_1,
        build_slide_2,
        build_slide_3,
        build_slide_4,
        build_slide_5,
        build_slide_6,
        build_slide_7,
        build_slide_8,
        build_slide_9,
        build_slide_10,
        build_slide_11,
        build_slide_12,
        build_slide_13,
        build_slide_14,
        build_slide_15,
    ]

    for i, builder in enumerate(builders, 1):
        print(f'  Building slide {i}/15 — {builder.__name__}...')
        builder(prs)

    out = '/home/user/yupebis/Avengers_Doomsday_Presentation.pptx'
    prs.save(out)
    print(f'\n✅  Saved → {out}')


if __name__ == '__main__':
    main()
