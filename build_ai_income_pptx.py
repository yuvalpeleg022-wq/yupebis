"""
AI Income Lab — PowerPoint Presentation Builder
Converts the AI Income Lab HTML site into a matching .pptx
Dark aesthetic: #0a0a0a bg · #c4ff3d lime accent · #f5f5f0 ink
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0A, 0x0A, 0x0A)
BG_CARD     = RGBColor(0x14, 0x14, 0x14)
BG_ELEV     = RGBColor(0x1A, 0x1A, 0x1A)
LINE        = RGBColor(0x2A, 0x2A, 0x2A)
INK         = RGBColor(0xF5, 0xF5, 0xF0)
INK_MUTE    = RGBColor(0xA8, 0xA8, 0x9E)
INK_FAINT   = RGBColor(0x5A, 0x5A, 0x52)
ACCENT      = RGBColor(0xC4, 0xFF, 0x3D)   # lime green
ACCENT_WARM = RGBColor(0xFF, 0x78, 0x49)   # orange
ACCENT_COOL = RGBColor(0x7A, 0xA7, 0xFF)   # blue
BLACK       = RGBColor(0x00, 0x00, 0x00)

# ── Slide size: widescreen 13.3 × 7.5 ────────────────────────────────────────
W = Inches(13.3)
H = Inches(7.5)


# ── XML helpers ───────────────────────────────────────────────────────────────

def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def _get_spPr(shape):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    return spPr


def solid_fill(shape, color: RGBColor, alpha_pct: int = 100):
    spPr = _get_spPr(shape)
    for tag in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'), qn('a:pattFill')]:
        e = spPr.find(tag)
        if e is not None:
            spPr.remove(e)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', rgb_hex(color))
    if alpha_pct < 100:
        a = etree.SubElement(sc, qn('a:alpha'))
        a.set('val', str(alpha_pct * 1000))


def no_fill(shape):
    spPr = _get_spPr(shape)
    for tag in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill')]:
        e = spPr.find(tag)
        if e is not None:
            spPr.remove(e)
    etree.SubElement(spPr, qn('a:noFill'))


def set_line(shape, color: RGBColor, width_pt: float = 1, alpha_pct: int = 100):
    spPr = _get_spPr(shape)
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(int(width_pt * 12700)))
    sf = etree.SubElement(ln, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', rgb_hex(color))
    if alpha_pct < 100:
        a = etree.SubElement(sc, qn('a:alpha'))
        a.set('val', str(alpha_pct * 1000))


def no_line(shape):
    spPr = _get_spPr(shape)
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:noFill'))


def grad_fill(shape, stops):
    """stops = [(pos_0_100, RGBColor, alpha_pct), ...]  — top-to-bottom"""
    spPr = _get_spPr(shape)
    for tag in [qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill')]:
        e = spPr.find(tag)
        if e is not None:
            spPr.remove(e)
    gf = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gf, qn('a:gsLst'))
    for pos, color, alpha_pct in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        sc = etree.SubElement(gs, qn('a:srgbClr'))
        sc.set('val', rgb_hex(color))
        if alpha_pct < 100:
            a = etree.SubElement(sc, qn('a:alpha'))
            a.set('val', str(alpha_pct * 1000))
    lin = etree.SubElement(gf, qn('a:lin'))
    lin.set('ang', '5400000')
    lin.set('scaled', '0')


def shadow(shape, blur_pt=18, dist_pt=4, angle_deg=135, opacity_pct=60):
    spPr = _get_spPr(shape)
    eLst = spPr.find(qn('a:effectLst'))
    if eLst is None:
        eLst = etree.SubElement(spPr, qn('a:effectLst'))
    os = etree.SubElement(eLst, qn('a:outerShdw'))
    os.set('blurRad', str(int(blur_pt * 12700)))
    os.set('dist', str(int(dist_pt * 12700)))
    os.set('dir', str(int(angle_deg * 60000)))
    os.set('algn', 'ctr')
    os.set('rotWithShape', '0')
    sc = etree.SubElement(os, qn('a:srgbClr'))
    sc.set('val', '000000')
    a = etree.SubElement(sc, qn('a:alpha'))
    a.set('val', str(opacity_pct * 1000))


def shape_glow(shape, color: RGBColor, radius_pt=12, alpha_pct=60):
    spPr = _get_spPr(shape)
    eLst = spPr.find(qn('a:effectLst'))
    if eLst is None:
        eLst = etree.SubElement(spPr, qn('a:effectLst'))
    gw = etree.SubElement(eLst, qn('a:glow'))
    gw.set('rad', str(int(radius_pt * 12700)))
    sc = etree.SubElement(gw, qn('a:srgbClr'))
    sc.set('val', rgb_hex(color))
    a = etree.SubElement(sc, qn('a:alpha'))
    a.set('val', str(alpha_pct * 1000))


def run_glow(run, color: RGBColor, radius_pt=8, alpha_pct=70):
    rPr = run._r
    eLst = rPr.find(qn('a:effectLst'))
    if eLst is None:
        eLst = etree.SubElement(rPr, qn('a:effectLst'))
    gw = etree.SubElement(eLst, qn('a:glow'))
    gw.set('rad', str(int(radius_pt * 12700)))
    sc = etree.SubElement(gw, qn('a:srgbClr'))
    sc.set('val', rgb_hex(color))
    a = etree.SubElement(sc, qn('a:alpha'))
    a.set('val', str(alpha_pct * 1000))


def rotate(shape, deg: float):
    spPr = _get_spPr(shape)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    xfrm.set('rot', str(int(deg * 60000)))


# ── Slide/shape factories ─────────────────────────────────────────────────────

RECT = 1
OVAL = 9

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def rect(slide, x, y, w, h):
    return slide.shapes.add_shape(RECT, x, y, w, h)


def oval(slide, x, y, w, h):
    return slide.shapes.add_shape(OVAL, x, y, w, h)


def card(slide, x, y, w, h, fill=BG_CARD, border=LINE, border_pt=1):
    s = rect(slide, x, y, w, h)
    solid_fill(s, fill)
    set_line(s, border, border_pt)
    shadow(s)
    return s


def accent_card(slide, x, y, w, h):
    """Card with lime-green left border accent."""
    s = card(slide, x, y, w, h, BG_ELEV, LINE)
    # left accent bar
    bar = rect(slide, x, y, Pt(3), h)
    solid_fill(bar, ACCENT)
    no_line(bar)
    return s


def hline(slide, x, y, w, color=LINE, pt=1):
    s = rect(slide, x, y, w, Pt(pt))
    solid_fill(s, color)
    no_line(s)
    return s


def label_tag(slide, x, y, w, text, color=ACCENT):
    """Small monospace section label."""
    tb = slide.shapes.add_textbox(x, y, w, Pt(22))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = 'Courier New'
    r.font.size = Pt(9)
    r.font.color.rgb = color
    r.font.bold = True
    hline(slide, x, y + Pt(24), w, color, 1)
    return tb


def title_run(slide, text, x, y, w, h, size=52, color=INK,
              font='Inter Tight', bold=True, align=PP_ALIGN.LEFT, glow=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    if glow:
        run_glow(r, color, 10, 60)
    return tb


def body_run(slide, text, x, y, w, h, size=13, color=INK_MUTE,
             font='Inter Tight', align=PP_ALIGN.LEFT, bold=False, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def multiline_body(slide, lines, x, y, w, h, size=13, color=INK_MUTE, spacing_pt=6):
    """Each item in lines is rendered as a separate paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(spacing_pt)
        r = p.add_run()
        r.text = line
        r.font.name = 'Inter Tight'
        r.font.size = Pt(size)
        r.font.color.rgb = color


def arrow_list(slide, items, x, y, w, h, size=12):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = '→  '
        r1.font.name = 'Courier New'
        r1.font.size = Pt(size)
        r1.font.color.rgb = ACCENT
        r2 = p.add_run()
        r2.text = item
        r2.font.name = 'Inter Tight'
        r2.font.size = Pt(size)
        r2.font.color.rgb = INK


def stat_block(slide, x, y, w, h, number, unit, desc, tag=''):
    c = card(slide, x, y, w, h, BG_CARD, LINE)
    # tag
    if tag:
        body_run(slide, tag, x + Inches(0.18), y + Pt(10), w - Inches(0.36), Pt(18),
                 size=9, color=INK_FAINT, font='Courier New')
    # number
    nb = slide.shapes.add_textbox(x + Inches(0.18), y + Pt(28), w - Inches(0.36), Inches(1.1))
    p = nb.text_frame.paragraphs[0]
    r1 = p.add_run()
    r1.text = number
    r1.font.name = 'Georgia'
    r1.font.size = Pt(58)
    r1.font.color.rgb = ACCENT
    r2 = p.add_run()
    r2.text = unit
    r2.font.name = 'Georgia'
    r2.font.size = Pt(28)
    r2.font.color.rgb = INK_MUTE
    # desc
    body_run(slide, desc, x + Inches(0.18), y + Inches(1.55), w - Inches(0.36), Pt(40),
             size=10, color=INK_MUTE, font='Courier New')


def method_card(slide, x, y, w, h, num, name, desc, level, pay, stack):
    c = card(slide, x, y, w, h, BG_CARD, LINE)
    body_run(slide, num, x + Inches(0.2), y + Pt(12), w - Inches(0.4), Pt(18),
             size=9, color=INK_FAINT, font='Courier New')
    title_run(slide, name, x + Inches(0.2), y + Pt(30), w - Inches(0.4), Pt(45),
              size=20, font='Georgia', bold=False)
    body_run(slide, desc, x + Inches(0.2), y + Pt(78), w - Inches(0.4), Pt(72),
             size=11.5, color=INK_MUTE)
    # chips
    lv = rect(slide, x + Inches(0.2), y + Inches(2.0), Inches(1.0), Pt(20))
    solid_fill(lv, RGBColor(0x1E, 0x28, 0x00), 100)
    set_line(lv, RGBColor(0x5A, 0x80, 0x00), 1)
    body_run(slide, level, x + Inches(0.22), y + Inches(2.01), Inches(0.96), Pt(18),
             size=9, color=ACCENT, font='Courier New')
    pv = rect(slide, x + Inches(1.28), y + Inches(2.0), Inches(1.4), Pt(20))
    solid_fill(pv, BG_ELEV)
    set_line(pv, LINE, 1)
    body_run(slide, pay, x + Inches(1.30), y + Inches(2.01), Inches(1.36), Pt(18),
             size=9, color=INK_MUTE, font='Courier New')
    # stack
    hline(slide, x + Inches(0.2), y + Inches(2.32), w - Inches(0.4), LINE, 1)
    body_run(slide, 'Stack — ' + stack, x + Inches(0.2), y + Inches(2.38), w - Inches(0.4), Pt(18),
             size=9, color=INK_FAINT, font='Courier New')


def tool_card(slide, x, y, w, h, category, name, purpose):
    c = card(slide, x, y, w, h, BG_CARD, LINE)
    body_run(slide, category, x + Inches(0.18), y + Pt(14), w - Inches(0.36), Pt(18),
             size=9, color=INK_FAINT, font='Courier New')
    title_run(slide, name, x + Inches(0.18), y + Pt(32), w - Inches(0.36), Pt(35),
              size=20, font='Georgia', bold=False)
    body_run(slide, purpose, x + Inches(0.18), y + Pt(70), w - Inches(0.36), Pt(38),
             size=11, color=INK_MUTE)


def phase_row(slide, y, h, num, week, title, desc, tasks):
    # phase number
    nb = slide.shapes.add_textbox(Inches(0.4), y, Inches(1.2), Inches(1.2))
    p = nb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = num
    r.font.name = 'Georgia'
    r.font.size = Pt(72)
    r.font.color.rgb = ACCENT
    body_run(slide, week, Inches(0.4), y + Inches(0.95), Inches(1.2), Pt(18),
             size=9, color=INK_FAINT, font='Courier New')
    # content
    title_run(slide, title, Inches(1.9), y + Pt(10), Inches(7.0), Pt(42),
              size=24, font='Georgia', bold=False)
    body_run(slide, desc, Inches(1.9), y + Pt(56), Inches(7.0), Pt(52),
             size=12, color=INK_MUTE)
    arrow_list(slide, tasks, Inches(1.9), y + Pt(112), Inches(6.8), Pt(60), size=11)
    hline(slide, Inches(0.4), y + h - Pt(2), Inches(12.5), LINE, 1)


def warn_card(slide, x, y, w, h, num, heading, body_text):
    # orange left bar
    bar = rect(slide, x, y, Pt(3), h)
    solid_fill(bar, ACCENT_WARM)
    no_line(bar)
    c = card(slide, x + Pt(3), y, w - Pt(3), h, BG_ELEV, LINE)
    body_run(slide, num, x + Pt(12), y + Pt(14), w - Pt(20), Pt(18),
             size=9, color=ACCENT_WARM, font='Courier New')
    title_run(slide, heading, x + Pt(12), y + Pt(30), w - Pt(20), Pt(38),
              size=20, font='Georgia', bold=False)
    body_run(slide, body_text, x + Pt(12), y + Pt(72), w - Pt(20), Pt(62),
             size=11.5, color=INK_MUTE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide"""
    s = new_slide(prs)

    # Subtle grid lines (decorative)
    for i in range(0, 14):
        vl = rect(s, Inches(i * 1.0), Inches(0), Pt(1), H)
        solid_fill(vl, LINE, 30)
        no_line(vl)
    for i in range(0, 8):
        hl = rect(s, Inches(0), Inches(i * 1.0), W, Pt(1))
        solid_fill(hl, LINE, 30)
        no_line(hl)

    # Logo box
    logo_bg = rect(s, Inches(0.5), Inches(0.5), Pt(36), Pt(36))
    solid_fill(logo_bg, ACCENT)
    no_line(logo_bg)
    lb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Pt(36), Pt(36))
    p = lb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'AI'
    r.font.name = 'Courier New'
    r.font.size = Pt(14)
    r.font.color.rgb = BLACK
    r.font.bold = True

    body_run(s, 'INCOME LAB / personal notebook',
             Inches(1.05), Inches(0.6), Inches(5), Pt(22),
             size=11, color=INK, font='Courier New', bold=True)

    # Ticker-style meta
    body_run(s, '// Field manual   ·   Edition 01 — May 2026',
             Inches(0.5), Inches(1.5), Inches(7), Pt(20),
             size=10, color=INK_FAINT, font='Courier New')

    # Live tag
    tag_bg = oval(s, Inches(0.5), Inches(2.05), Inches(3.0), Pt(26))
    solid_fill(tag_bg, RGBColor(0x14, 0x28, 0x00), 90)
    set_line(tag_bg, ACCENT, 1)
    body_run(s, '● Live transmission — AI economy 2026',
             Inches(0.62), Inches(2.07), Inches(2.8), Pt(22),
             size=9.5, color=ACCENT, font='Courier New')

    # Main headline
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(10), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = 'Making money online,\npowered by '
    r1.font.name = 'Georgia'
    r1.font.size = Pt(62)
    r1.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = 'artificial intelligence'
    r2.font.name = 'Georgia'
    r2.font.size = Pt(62)
    r2.font.color.rgb = ACCENT
    r2.font.italic = True

    # Subtitle
    body_run(s,
        'A working notebook of the methods, platforms, and tools that actually generate income '
        'in 2026 — researched, tested, and documented for personal study. No hype. Just leverage.',
        Inches(0.5), Inches(5.85), Inches(9.0), Pt(55),
        size=15, color=INK_MUTE)

    # Right-side edition box
    eb = card(s, Inches(10.5), Inches(2.8), Inches(2.4), Inches(2.0), BG_ELEV, LINE)
    body_run(s, '// operator', Inches(10.65), Inches(2.95), Inches(2.1), Pt(18),
             size=9, color=INK_FAINT, font='Courier New')
    body_run(s, 'Personal\nstudy log', Inches(10.65), Inches(3.2), Inches(2.1), Inches(0.8),
             size=18, color=INK, font='Georgia')


def slide_02_market(prs):
    """01 / Market — 4 stat blocks"""
    s = new_slide(prs)

    label_tag(s, Inches(0.5), Inches(0.45), Inches(2), '01 / Market')

    tb = s.shapes.add_textbox(Inches(2.5), Inches(0.3), Inches(9.5), Inches(1.1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = 'The opportunity, '
    r1.font.name = 'Georgia'
    r1.font.size = Pt(40)
    r1.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = 'by the numbers'
    r2.font.name = 'Georgia'
    r2.font.size = Pt(40)
    r2.font.color.rgb = INK_MUTE
    r2.font.italic = True

    body_run(s,
        'Why now matters. The gap between tools and the people who know how to wield them '
        'is the largest it has ever been — and businesses are paying to close it.',
        Inches(2.5), Inches(1.45), Inches(9.5), Pt(44),
        size=14, color=INK_MUTE)

    hline(s, Inches(0.5), Inches(2.3), Inches(12.3), LINE)

    stats = [
        ('674', 'B',  'Size of the global gig economy in 2026',  '$ — global'),
        ('70',  '%',  'Of firms increasing AI spending this year', '% — companies'),
        ('15.7','T',  "AI's projected contribution to GDP by 2030", '$ — projection'),
        ('300', '%',  'Increase in AI job listings since 2023',   '% — growth'),
    ]
    sw = Inches(3.0)
    gap = Inches(0.12)
    sy = Inches(2.5)
    for i, (num, unit, desc, tag) in enumerate(stats):
        stat_block(s, Inches(0.4) + i * (sw + gap), sy, sw, Inches(2.2), num, unit, desc, tag)

    # Bottom intro rows
    body_run(s,
        'AI consulting rates: $100–$300 / hr   ·   Prompt engineering: $50–$150 / hr   ·   '
        'Most digital products on Gumroad take days, not weeks',
        Inches(0.5), Inches(6.85), Inches(12.3), Pt(22),
        size=10, color=INK_FAINT, font='Courier New', align=PP_ALIGN.CENTER)


def slide_03_methods_overview(prs):
    """02 / Methods — overview header + first 5"""
    s = new_slide(prs)

    label_tag(s, Inches(0.5), Inches(0.3), Inches(2.2), '02 / Methods')

    tb = s.shapes.add_textbox(Inches(2.8), Inches(0.15), Inches(9.8), Inches(1.1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'Ten ways to earn, '
    r1.font.name = 'Georgia'; r1.font.size = Pt(36); r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = 'ranked by leverage'
    r2.font.name = 'Georgia'; r2.font.size = Pt(36)
    r2.font.color.rgb = INK_MUTE; r2.font.italic = True

    body_run(s,
        'Every method pairs a real human skill with an AI tool that multiplies output. '
        'The goal isn\'t to be a "prompt engineer" — it\'s to deliver outcomes businesses already pay for, faster.',
        Inches(2.8), Inches(1.3), Inches(9.8), Pt(38), size=12.5, color=INK_MUTE)

    hline(s, Inches(0.5), Inches(2.1), Inches(12.3), LINE)

    methods_a = [
        ('001 — content & writing', 'AI-powered content services',
         'Ghost-writing blogs, newsletters, and landing copy for businesses. AI handles research and first drafts. You handle voice, strategy, and quality control.',
         'Beginner', '$40–$100 / hr', 'ChatGPT · Surfer SEO · Grammarly'),
        ('002 — visual design', 'AI art & print-on-demand',
         'Generate distinctive visuals with Midjourney, then sell them as prints, T-shirts, or downloadable art on Etsy, Redbubble, and Printify. Pick a narrow niche.',
         'Beginner', '$500–$3K / mo', 'Midjourney · DALL·E · Printify · Etsy'),
        ('003 — digital products', 'Templates, prompt packs & eBooks',
         'Build once, sell forever. Notion templates, prompt libraries, spreadsheet trackers, resume kits. The clearest path to genuinely semi-passive income.',
         'Beginner', 'Semi-passive', 'Notion · Gumroad · Canva · Etsy'),
        ('004 — video', 'Faceless YouTube channels',
         'Scripted with AI, narrated with ElevenLabs, edited with Pictory or Descript. Best for educational, list-style, or storytelling content. Long ramp, compounds powerfully.',
         'Intermediate', '$150–$2K / mo', 'ChatGPT · ElevenLabs · Pictory · CapCut'),
        ('005 — voice', 'Voiceover & localization',
         'Sell narration, dubbing, and translation for short-form video, audiobooks, and ads. High-margin work that AI completes in minutes.',
         'Intermediate', '$50–$500 / project', 'ElevenLabs · Murf · Descript'),
    ]

    mw = Inches(2.54)
    gap = Inches(0.07)
    mh = Inches(2.75)
    for i, (num, name, desc, lvl, pay, stack) in enumerate(methods_a):
        method_card(s, Inches(0.4) + i * (mw + gap), Inches(2.2),
                    mw, mh, num, name, desc, lvl, pay, stack)


def slide_04_methods_b(prs):
    """02 / Methods — methods 6-10"""
    s = new_slide(prs)

    label_tag(s, Inches(0.5), Inches(0.3), Inches(2.2), '02 / Methods (cont.)')
    body_run(s, 'Advanced & high-leverage income methods', Inches(2.8), Inches(0.3),
             Inches(9.8), Pt(30), size=32, color=INK, font='Georgia')
    hline(s, Inches(0.5), Inches(1.05), Inches(12.3), LINE)

    methods_b = [
        ('006 — marketing', 'Affiliate marketing with AI SEO',
         'Build niche review sites or YouTube channels. Use AI to research keywords, draft content, and generate B-roll. Best after you have your own audience or media asset.',
         'Intermediate', '20–50% commission', 'Surfer SEO · Ahrefs · ChatGPT · WordPress'),
        ('007 — education', 'Online courses & cohorts',
         'Teach one skill exceptionally well. AI compresses curriculum design, slide creation, and marketing copy. The highest-margin product in this entire list.',
         'Intermediate', '$1K–$10K / course', 'Thinkific · Coursebox · Loom · Notion'),
        ('008 — agents', 'Custom GPTs & AI agents',
         'Build narrow-purpose AI agents for small businesses — customer support bots, lead qualifiers, internal knowledge bases. Sell to a single industry to compound expertise.',
         'Advanced', '$75–$200 / hr', 'OpenAI Assistants · Voiceflow'),
        ('009 — automation', 'Business workflow automation',
         'Wire AI into operations. Email triage, meeting notes, lead routing, content pipelines. Service businesses pay 4-figure retainers for systems that save them hours weekly.',
         'Advanced', '$75–$200 / hr', 'Make · n8n · Zapier · Airtable'),
        ('010 — consulting', 'AI consulting & training',
         'Help organizations identify where AI creates real ROI. Run workshops. Audit their workflows. The highest hourly rate in the entire field — built on trust and credibility.',
         'Advanced', '$100–$300 / hr', 'Your expertise + Loom + slides'),
    ]

    mw = Inches(2.54)
    gap = Inches(0.07)
    mh = Inches(2.75)
    for i, (num, name, desc, lvl, pay, stack) in enumerate(methods_b):
        method_card(s, Inches(0.4) + i * (mw + gap), Inches(1.25),
                    mw, mh, num, name, desc, lvl, pay, stack)

    body_run(s,
        'Start at 001 or 002. Run that single method for 60 days before evaluating. Switching early is the #1 mistake.',
        Inches(0.5), Inches(4.2), Inches(12.3), Pt(28),
        size=13, color=INK_MUTE, align=PP_ALIGN.CENTER)


def slide_05_toolkit(prs):
    """03 / Toolkit"""
    s = new_slide(prs)

    label_tag(s, Inches(0.5), Inches(0.3), Inches(2), '03 / Toolkit')

    tb = s.shapes.add_textbox(Inches(2.8), Inches(0.15), Inches(9.8), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'The working stack, '
    r1.font.name = 'Georgia'; r1.font.size = Pt(36); r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = 'nothing more'
    r2.font.name = 'Georgia'; r2.font.size = Pt(36)
    r2.font.color.rgb = INK_MUTE; r2.font.italic = True

    body_run(s,
        'Most operators run on fewer than five tools. Mastery over a small set beats shallow familiarity with twenty. Start with the first four.',
        Inches(2.8), Inches(1.1), Inches(9.8), Pt(28), size=13, color=INK_MUTE)

    hline(s, Inches(0.5), Inches(1.75), Inches(12.3), LINE)

    tools = [
        ('Language model', 'ChatGPT',        'Daily driver for writing, research, and ideation.'),
        ('Language model', 'Gemini',          'Long documents, multimodal input, analysis.'),
        ('Image gen',      'Midjourney',      'High-end visuals, brand imagery, product art.'),
        ('Image gen',      'DALL·E',          'Quick illustrations, mockups, social content.'),
        ('Voice',          'ElevenLabs',      'Realistic narration in any language.'),
        ('Video',          'Pictory',         'Turn articles into short-form video, fast.'),
        ('Design',         'Canva AI',        'Templates, social posts, simple branding.'),
        ('Workspace',      'Notion AI',       'Your second brain — research, drafts, planning.'),
        ('Automation',     'Make',            'Visual workflows linking AI into your apps.'),
        ('Storefront',     'Gumroad',         'Sell digital products in under an hour.'),
        ('Marketplace',    'Etsy',            'Templates, printables, and AI art with traffic.'),
        ('Freelance',      'Upwork & Fiverr', 'First five clients without an audience.'),
    ]

    tw = Inches(3.0)
    th = Inches(1.1)
    gap_x = Inches(0.12)
    gap_y = Inches(0.1)
    sx, sy = Inches(0.4), Inches(1.9)
    for i, (cat, name, purpose) in enumerate(tools):
        col = i % 4
        row = i // 4
        x = sx + col * (tw + gap_x)
        y = sy + row * (th + gap_y)
        tool_card(s, x, y, tw, th, cat, name, purpose)


def slide_06_roadmap(prs):
    """04 / Roadmap — 60-day plan"""
    s = new_slide(prs)

    label_tag(s, Inches(0.5), Inches(0.3), Inches(2.2), '04 / Roadmap')

    tb = s.shapes.add_textbox(Inches(2.8), Inches(0.15), Inches(9.8), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'A 60-day field plan, '
    r1.font.name = 'Georgia'; r1.font.size = Pt(32); r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = 'step by step'
    r2.font.name = 'Georgia'; r2.font.size = Pt(32)
    r2.font.color.rgb = INK_MUTE; r2.font.italic = True

    body_run(s,
        'Pick ONE method and run this protocol for two months before evaluating. '
        'Switching methods early is the single biggest mistake new operators make.',
        Inches(2.8), Inches(1.05), Inches(9.8), Pt(28), size=12.5, color=INK_MUTE)

    hline(s, Inches(0.5), Inches(1.75), Inches(12.3), LINE)

    phases = [
        ('01', 'Week 1',   'Pick a single lane and commit',
         'Match your existing skill with one method. Writing? → 001 or 007. Visual? → 002 or 003. Tech? → 008 or 009.',
         ['List your top 3 skills', 'Cross-reference with the 10 methods', 'Pick one. Close the other tabs.']),
        ('02', 'Wk 1–2',  'Master two tools, deeply',
         'Deep fluency with two tools beats shallow knowledge of ten. Thirty minutes a day for 14 days. Build a personal prompt library.',
         ['Pick your primary LLM', 'Pick one secondary tool', 'Save every prompt that produced a good result']),
        ('03', 'Wk 2–3',  'Ship a small first deliverable',
         'A prompt pack. A Notion template. Five sample articles. Proof to yourself and the market that you can deliver.',
         ['Define the smallest possible deliverable', 'Build it in under a week', 'Publish it somewhere public']),
        ('04', 'Wk 3–4',  'First three clients at a low rate',
         'Trade price for testimonials. Three real clients with written feedback unlocks your real pricing tier. Don\'t skip this.',
         ['Set up Upwork / Fiverr / LinkedIn', 'Send 10 personalized pitches per day', 'Collect a testimonial after every job']),
        ('05', 'Wk 5–6',  'Daily outreach + daily content',
         'Inbound deal flow comes from being visible. One platform, one post format, daily. Same for outreach — ten conversations a day.',
         ['Pick one platform: LinkedIn or X', 'One post per day, same format', 'Ten outbound conversations daily']),
        ('06', 'Wk 6–7',  'Sell outcomes, not hours',
         '"A blog that brings 10 qualified leads a month" is worth more than "a 1,000-word article." Reframe every offer around the business result.',
         ['Rewrite your offer as an outcome', 'Productize at three price points', 'Test the new framing on five prospects']),
        ('07', 'Wk 7–8',  'Layer automation into operations',
         'Once delivery is repeatable, wire it together. Make or Zapier. Pre-built prompt chains. This step lets you double income without doubling time.',
         ['Map your service as a flowchart', 'Identify the three slowest steps', 'Automate or templatize each one']),
    ]

    ph = Inches(0.73)
    for i, (num, week, title, desc, tasks) in enumerate(phases):
        y = Inches(1.88) + i * ph
        phase_row(s, y, ph, num, week, title, desc, tasks)


def slide_07_caveats(prs):
    """05 / Caveats"""
    s = new_slide(prs)

    label_tag(s, Inches(0.5), Inches(0.3), Inches(2.2), '05 / Caveats')

    tb = s.shapes.add_textbox(Inches(2.8), Inches(0.15), Inches(9.8), Inches(1.1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'What no one selling courses '
    r1.font.name = 'Georgia'; r1.font.size = Pt(34); r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = 'will tell you'
    r2.font.name = 'Georgia'; r2.font.size = Pt(34)
    r2.font.color.rgb = INK_MUTE; r2.font.italic = True

    body_run(s,
        'Real friction lives in places the influencers skip past. Internalize these before you spend a dollar on a "passive income" program.',
        Inches(2.8), Inches(1.3), Inches(9.8), Pt(36), size=13, color=INK_MUTE)

    hline(s, Inches(0.5), Inches(2.1), Inches(12.3), LINE)

    warns = [
        ('Caveat 01', '"Easy money" is the red flag',
         'Anyone promising $10K/month passive without effort is selling you a course about earning $10K/month passive. Real businesses take 3 to 6 months of focused work before meaningful revenue.'),
        ('Caveat 02', 'Fact-check every AI output',
         'Generative models confabulate. Numbers, quotes, and citations must be verified before they reach a client. The single fastest way to destroy a reputation is one hallucinated statistic.'),
        ('Caveat 03', 'The market is flooded with slop',
         'Generic AI output is already a commodity. Winning means narrow niche + deep insight + production quality — not volume. Compete on taste and judgment, not throughput.'),
        ('Caveat 04', 'Read the licensing terms',
         'Commercial-use rights on AI-generated content shift frequently across tools and jurisdictions. Verify your right to sell before you list a single product.'),
    ]

    ww = Inches(6.1)
    wh = Inches(2.25)
    for i, (num, heading, body_text) in enumerate(warns):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * (ww + Inches(0.2))
        y = Inches(2.2) + row * (wh + Inches(0.15))
        warn_card(s, x, y, ww, wh, num, heading, body_text)


def slide_08_closing(prs):
    """Closing slide"""
    s = new_slide(prs)

    # Grid decoration
    for i in range(0, 14):
        vl = rect(s, Inches(i * 1.0), Inches(0), Pt(1), H)
        solid_fill(vl, LINE, 20)
        no_line(vl)

    # Big closing statement
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = 'Start small. Ship daily. '
    r1.font.name = 'Georgia'; r1.font.size = Pt(52); r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = 'Compound monthly.'
    r2.font.name = 'Georgia'; r2.font.size = Pt(52)
    r2.font.color.rgb = ACCENT; r2.font.italic = True

    body_run(s,
        'This notebook is a living document built to help think clearly about how to use AI to '
        'create real income. Updated as the landscape shifts.',
        Inches(0.8), Inches(4.6), Inches(9.0), Pt(44),
        size=15, color=INK_MUTE)

    hline(s, Inches(0.8), Inches(5.7), Inches(11.7), LINE)

    # Sections summary
    sections = ['01 Market', '02 Methods', '03 Toolkit', '04 Roadmap', '05 Caveats']
    for i, sec in enumerate(sections):
        body_run(s, sec, Inches(0.8) + i * Inches(2.4), Inches(5.8),
                 Inches(2.2), Pt(22), size=10, color=INK_FAINT, font='Courier New')

    body_run(s, '// AI Income Lab — Personal study site   ·   Built May 2026 · v0.2',
             Inches(0.8), Inches(7.1), Inches(11.7), Pt(20),
             size=9, color=INK_FAINT, font='Courier New', align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    builders = [
        ('Title',            slide_01_title),
        ('Market / Stats',   slide_02_market),
        ('Methods 1–5',      slide_03_methods_overview),
        ('Methods 6–10',     slide_04_methods_b),
        ('Toolkit',          slide_05_toolkit),
        ('Roadmap',          slide_06_roadmap),
        ('Caveats',          slide_07_caveats),
        ('Closing',          slide_08_closing),
    ]

    for i, (name, fn) in enumerate(builders, 1):
        print(f'  [{i}/{len(builders)}] {name}...')
        fn(prs)

    out = '/home/user/yupebis/AI_Income_Lab_Presentation.pptx'
    prs.save(out)
    print(f'\n✅  Saved → {out}')


if __name__ == '__main__':
    main()
